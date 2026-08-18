# -*- coding: utf-8 -*-
"""Genera 'Resumen Reportes - Matriz Mayo Julio 2026.pptx' con la matriz de
reportes (reportes que se mandan.xlsx): resumen, matriz completa, peso de
tareas por responsable y distribución."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "Resumen Reportes - Matriz Mayo Julio 2026.pptx"

# ---------- paleta ----------
DARK = RGBColor(0x1F, 0x4E, 0x79)
MID = RGBColor(0x2E, 0x75, 0xB6)
LIGHT = RGBColor(0xDE, 0xEB, 0xF7)
TEXT = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x80, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CAT_COLOR = {
    "Cobros y Pagos": RGBColor(0xF4, 0xB1, 0x83),
    "Sinergia y Renovación": RGBColor(0xA9, 0xD1, 0x8E),
    "Operaciones SOHO": RGBColor(0x9D, 0xC3, 0xE6),
    "Ventas / Cierre": RGBColor(0xFF, 0xE6, 0x99),
    "Monitoreo/Otros": RGBColor(0xC9, 0xA0, 0xDC),
}

REPORTS = [
    ("Cobros y Pagos", "Promesas de Pago", "Diario", "09:00", "Leonardo Sanabria"),
    ("Cobros y Pagos", "Reporte de Cobros", "2x Semana (Mar/Vie)", "15:00", "Edgar S. Centella / Luis Figueroa"),
    ("Sinergia y Renovación", "Sinergia, P2P y Renovación", "Diario", "Toda la jornada", "Raziel Nogal / Rosa Martínez"),
    ("Operaciones SOHO", "Reporte Llamadas (Diario)", "Diario", "12:00 - 13:00", "Leonardo Sanabria"),
    ("Operaciones SOHO", "Reporte Llamadas (Mensual)", "Diario", "09:00 - 10:00", "Leonardo Sanabria"),
    ("Ventas / Cierre", "Cierre Ventas Dataphone", "Diario", "18:00 - 19:00", "Leonardo Sanabria"),
    ("Monitoreo/Otros", "Sharep", "Diario", "Antes de 09:00 AM", "Raziel Nogal"),
    ("Monitoreo/Otros", "Balance día anterior", "Diario", "20:00 - 22:00", "Leonardo Sanabria"),
    ("Monitoreo/Otros", "Asesores Fuera Cumplimiento", "Diario", "20:00 - 22:00", "Leonardo Sanabria"),
    ("Monitoreo/Otros", "Ranking día previo", "Diario", "20:00 - 22:00", "Leonardo Sanabria"),
]

RESPONSABLES = [
    ("Leonardo Sanabria", 7, RGBColor(0x1F, 0x4E, 0x79),
     "Promesas de Pago · Reporte Llamadas (diario y mensual) · Cierre Ventas Dataphone · Balance día anterior · Asesores Fuera Cumplimiento · Ranking día previo"),
    ("Raziel Nogal", 2, RGBColor(0x54, 0x82, 0x35),
     "Sharep · Sinergia, P2P y Renovación (con Rosa Martínez)"),
    ("Edgar S. Centella / Luis Figueroa", 1, RGBColor(0xBF, 0x8F, 0x00),
     "Reporte de Cobros"),
]

CATS = [
    ("Monitoreo/Otros", 4),
    ("Cobros y Pagos", 2),
    ("Operaciones SOHO", 2),
    ("Ventas / Cierre", 1),
    ("Sinergia y Renovación", 1),
]


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, fill, line=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def bar(slide, x, y, w, h, fill, pct=None):
    """Barra simple (rectángulo redondeado)."""
    return rect(slide, x, y, w, h, fill, radius=0.25)


def set_cell(cell, text, size=10, bold=False, color=TEXT, fill=None, align=PP_ALIGN.LEFT):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.margin_left = cell.margin_right = Inches(0.08)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, titulo, sub=None):
    textbox(slide, 0.55, 0.32, 12.2, 0.7, titulo, size=30, bold=True, color=DARK)
    rect(slide, 0.6, 1.02, 2.4, 0.055, MID, radius=0.5)
    if sub:
        textbox(slide, 0.55, 1.15, 12.2, 0.4, sub, size=13, color=GRAY)


def footer(slide, texto="Gerencia de Operaciones · Informe de entregables"):
    textbox(slide, 0.55, 7.05, 12.2, 0.3, texto, size=9, color=GRAY)


# =====================================================================
def slide_titulo(prs):
    s = blank(prs)
    rect(s, 0, 0, 13.333, 0.28, DARK, radius=0)
    rect(s, 0, 7.22, 13.333, 0.28, DARK, radius=0)
    textbox(s, 1.0, 2.3, 11.3, 1.4, "Consolidado de Reportes", size=48, bold=True, color=DARK)
    textbox(s, 1.0, 3.6, 11.3, 0.7, "Matriz de reportes que se envían · Mayo – Julio 2026", size=20, color=MID)
    rect(s, 1.05, 4.55, 2.6, 0.06, MID, radius=0.5)
    textbox(s, 1.0, 4.85, 11.3, 0.5, "Gerencia de Operaciones", size=14, color=GRAY)


# =====================================================================
def slide_resumen(prs):
    s = blank(prs)
    add_title(s, "Resumen ejecutivo", "Qué se entrega, con qué frecuencia y quién lo hace")
    kpis = [
        ("10", "Reportes en la matriz"),
        ("5", "Categorías de reportes"),
        ("9", "Diarios"),
        ("1", "Bisemanal (Mar/Vie)"),
        ("3", "Responsables / equipos"),
        ("20:00–22:00", "Horario pico de envío"),
    ]
    x = 0.5
    for num, label in kpis:
        card = rect(s, x, 1.75, 1.95, 1.65, LIGHT, radius=0.1)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = DARK
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT
        x += 2.07

    textbox(s, 0.55, 3.85, 12.2, 0.4, "¿Qué cubre?", size=16, bold=True, color=DARK)
    lineas = [
        "•  Cobros y Pagos: promesas de pago y reporte de cobros (2x semana).",
        "•  Sinergia y Renovación: seguimiento Sinergia, P2P y renovaciones durante toda la jornada.",
        "•  Operaciones SOHO: reporte de llamadas diario y mensual.",
        "•  Ventas / Cierre: cierre de ventas Dataphone al final del día.",
        "•  Monitoreo/Otros: Sharep, balance del día anterior, asesores fuera de cumplimiento y ranking del día previo (bloque nocturno 20:00–22:00).",
        "",
        "Fuente: matriz “reportes que se mandan.xlsx” (Consolidado Mayo – Julio 2026).",
    ]
    y = 4.35
    for linea in lineas:
        textbox(s, 0.65, y, 12.1, 0.35, linea, size=13, color=(TEXT if linea else GRAY))
        y += 0.36
    footer(s)


# =====================================================================
def slide_matriz(prs):
    s = blank(prs)
    add_title(s, "Matriz de reportes", "Consolidado de Reportes (Mayo – Julio 2026) · 10 reportes")

    rows, cols = 11, 5
    tbl = s.shapes.add_table(rows, cols, Inches(0.45), Inches(1.75), Inches(12.45), Inches(4.9)).table
    tbl.first_row = False
    tbl.horz_banding = False

    widths = [2.5, 3.7, 1.75, 1.9, 2.6]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    tbl.rows[0].height = Inches(0.45)
    for i in range(1, rows):
        tbl.rows[i].height = Inches(0.44)

    headers = ["Categoría", "Nombre del Reporte", "Frecuencia", "Hora aprox.", "Responsable de Envío"]
    for c, h in enumerate(headers):
        set_cell(tbl.cell(0, c), h, size=11, bold=True, color=WHITE, fill=DARK, align=PP_ALIGN.CENTER)

    for i, (cat, nombre, freq, hora, resp) in enumerate(REPORTS, start=1):
        set_cell(tbl.cell(i, 1), nombre, size=10, bold=True)
        set_cell(tbl.cell(i, 2), freq, size=10)
        set_cell(tbl.cell(i, 3), hora, size=10)
        set_cell(tbl.cell(i, 4), resp, size=10)

    # agrupar categorías (filas contiguas) y colorear
    filas_cat = {}
    for i, (cat, *_rest) in enumerate(REPORTS, start=1):
        filas_cat.setdefault(cat, []).append(i)
    for cat, filas in filas_cat.items():
        first, last = filas[0], filas[-1]
        if last > first:
            tbl.cell(first, 0).merge(tbl.cell(last, 0))
        color = CAT_COLOR[cat]
        set_cell(tbl.cell(first, 0), cat, size=10, bold=True, color=TEXT, fill=color,
                 align=PP_ALIGN.CENTER)
    footer(s)


# =====================================================================
def slide_responsables(prs):
    s = blank(prs)
    add_title(s, "Peso de tareas por responsable", "Qué hace cada uno de los chicos y cuánto pesa en la operación")

    max_n = max(r[1] for r in RESPONSABLES)
    bar_ancho_max = 6.2

    # panel izquierdo: barras
    textbox(s, 0.55, 1.75, 6.5, 0.4, "Reportes por responsable", size=15, bold=True, color=DARK)
    y = 2.25
    for nombre, n, color, _det in RESPONSABLES:
        textbox(s, 0.55, y, 3.4, 0.35, nombre, size=12, bold=True, color=TEXT)
        bar(s, 3.95, y + 0.02, bar_ancho_max * n / max_n, 0.3, color)
        textbox(s, 3.95 + bar_ancho_max * n / max_n + 0.1, y - 0.04, 1.2, 0.4,
                f"{n} ({round(n * 100 / sum(r[1] for r in RESPONSABLES))}%)", size=13, bold=True, color=DARK)
        y += 0.62

    textbox(s, 0.55, y + 0.15, 6.6, 0.8,
            f"{sum(r[1] for r in RESPONSABLES)} reportes en total · 9 de 10 se envían a diario.",
            size=12, color=GRAY)

    # panel derecho: detalle por responsable
    textbox(s, 7.35, 1.75, 5.5, 0.4, "Reportes a cargo", size=15, bold=True, color=DARK)
    y = 2.25
    for nombre, n, color, det in RESPONSABLES:
        card = rect(s, 7.35, y, 5.5, 1.5, LIGHT, radius=0.08)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{nombre}  ·  {n} reporte{'s' if n != 1 else ''}"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = det
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT
        y += 1.65
    footer(s)


# =====================================================================
def slide_distribucion(prs):
    s = blank(prs)
    add_title(s, "Distribución de reportes", "Por categoría y por frecuencia")

    # izquierda: por categoría (barras horizontales)
    textbox(s, 0.55, 1.75, 6.5, 0.4, "Por categoría", size=15, bold=True, color=DARK)
    max_n = max(n for _c, n in CATS)
    y = 2.3
    for cat, n in CATS:
        textbox(s, 0.55, y, 3.0, 0.35, cat, size=12, color=TEXT)
        bar(s, 3.6, y + 0.02, 4.0 * n / max_n, 0.32, CAT_COLOR[cat])
        textbox(s, 3.6 + 4.0 * n / max_n + 0.1, y - 0.04, 0.9, 0.4, str(n), size=13, bold=True, color=DARK)
        y += 0.62

    # derecha: frecuencia + horario pico
    textbox(s, 8.0, 1.75, 4.8, 0.4, "Por frecuencia", size=15, bold=True, color=DARK)
    y = 2.3
    for label, n, color in [("Diarios", 9, MID), ("Bisemanales (Mar/Vie)", 1, RGBColor(0xBF, 0x8F, 0x00))]:
        card = rect(s, 8.0, y, 4.8, 1.2, LIGHT, radius=0.1)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(n)
        r.font.size = Pt(32)
        r.font.bold = True
        r.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(12)
        r2.font.color.rgb = TEXT
        y += 1.4

    textbox(s, 8.0, y + 0.15, 4.8, 0.9,
            "Horario pico de envío: 20:00 – 22:00\n(3 reportes nocturnos: balance, cumplimiento y ranking)",
            size=12, color=GRAY)
    footer(s)


# =====================================================================
def slide_cierre(prs):
    s = blank(prs)
    rect(s, 0, 0, 13.333, 0.28, DARK, radius=0)
    rect(s, 0, 7.22, 13.333, 0.28, DARK, radius=0)
    textbox(s, 1.0, 2.7, 11.3, 1.0, "Gracias", size=48, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    textbox(s, 1.0, 3.9, 11.3, 0.6,
            "10 reportes · 9 diarios · 3 responsables · mayor carga operativa: Leonardo Sanabria (7 reportes)",
            size=15, color=GRAY, align=PP_ALIGN.CENTER)


def main():
    prs = new_prs()
    slide_titulo(prs)
    slide_resumen(prs)
    slide_matriz(prs)
    slide_responsables(prs)
    slide_distribucion(prs)
    slide_cierre(prs)
    prs.save(OUT)
    print(f"OK -> {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas)")


if __name__ == "__main__":
    main()
